"""
make_slides.py — 20-slide main deck + 26 appendix slides (A-Z).
Run: python3 make_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
TEAL   = RGBColor(0x00, 0x88, 0x8A)
GOLD   = RGBColor(0xF0, 0xAB, 0x00)
GREEN  = RGBColor(0x27, 0x96, 0x58)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xD6, 0xE4, 0xF0)
DARK   = RGBColor(0x22, 0x22, 0x33)
GREY   = RGBColor(0x88, 0x88, 0x99)
SILVER = RGBColor(0xEC, 0xF0, 0xF1)
INK    = RGBColor(0x10, 0x1E, 0x33)
PURPLE = RGBColor(0x6C, 0x35, 0x8A)

SW, SH = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════════
# Core helpers
# ══════════════════════════════════════════════════════════════════════════════

def R(sl, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def T(sl, text, x, y, w, h, sz=18, bold=False, italic=False,
      color=DARK, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return t

def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text

def footer(sl):
    R(sl, 0, SH - Inches(0.25), SW, Inches(0.25), fill=NAVY)
    T(sl, "Columbia University  ·  Digital Marketplaces  ·  Spring 2026",
      Inches(0.3), SH - Inches(0.23), Inches(12.7), Inches(0.22),
      sz=9, color=RGBColor(0x7F, 0xA7, 0xCC), align=PP_ALIGN.LEFT)

def hdr(sl, text):
    R(sl, 0, 0, SW, Inches(0.95), fill=NAVY)
    R(sl, 0, Inches(0.95), SW, Pt(4), fill=GOLD)
    T(sl, text, Inches(0.4), Inches(0.1), Inches(12.1), Inches(0.8),
      sz=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def set_bar_colors(chart, colors_per_series):
    plot = chart.plots[0]
    for si, series in enumerate(plot.series):
        ser_el = series._element
        if si >= len(colors_per_series):
            continue
        for pi, color in enumerate(colors_per_series[si]):
            dPt = etree.SubElement(ser_el, qn('c:dPt'))
            idx = etree.SubElement(dPt, qn('c:idx'))
            idx.set('val', str(pi))
            b3d = etree.SubElement(dPt, qn('c:bubble3D'))
            b3d.set('val', '0')
            spPr = etree.SubElement(dPt, qn('c:spPr'))
            sf = etree.SubElement(etree.SubElement(spPr, qn('a:solidFill')), qn('a:srgbClr'))
            sf.set('val', str(color))
            ln = etree.SubElement(spPr, qn('a:ln'))
            etree.SubElement(etree.SubElement(ln, qn('a:solidFill')), qn('a:srgbClr')).set('val', str(color))

def add_chart(sl, chart_type, data, x, y, w, h):
    return sl.shapes.add_chart(chart_type, x, y, w, h, data).chart

def th(sl, cols, x, y, widths, h=Inches(0.4)):
    cx = x
    for col, w in zip(cols, widths):
        R(sl, cx, y, w, h, fill=NAVY)
        T(sl, col, cx + Inches(0.06), y + Inches(0.04),
          w - Inches(0.12), h - Inches(0.08),
          sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += w

def tr(sl, vals, x, y, widths, aligns=None, h=Inches(0.42), bg=WHITE, sz=11):
    if aligns is None:
        aligns = [PP_ALIGN.CENTER] * len(vals)
    cx = x
    for val, w, al in zip(vals, widths, aligns):
        R(sl, cx, y, w, h, fill=bg, line=GREY, lw=Pt(0.5))
        T(sl, val, cx + Inches(0.06), y + Inches(0.04),
          w - Inches(0.12), h - Inches(0.08),
          sz=sz, color=DARK, align=al)
        cx += w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
R(sl, 0, Inches(3.0), SW, Pt(5), fill=GOLD)

T(sl, "AI-Mediated Food Delivery", Inches(1), Inches(0.7), Inches(11.3), Inches(1.3),
  sz=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "Marketplaces", Inches(1), Inches(1.9), Inches(11.3), Inches(1.0),
  sz=46, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(sl, "Monetization Design  ·  Sponsorship Bias  ·  Consumer Welfare",
  Inches(1), Inches(3.2), Inches(11.3), Inches(0.55),
  sz=18, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

# Discovery flow strip
R(sl, Inches(1.0), Inches(4.1), Inches(11.33), Inches(1.2), fill=INK)
for i, (label, sub) in enumerate([
    ("Consumer", "states intent"),
    ("LLM", "discovery layer"),
    ("Platform", "fulfills order"),
    ("Restaurant", "earns revenue"),
]):
    bx = Inches(1.2) + i * Inches(2.75)
    R(sl, bx, Inches(4.2), Inches(2.4), Inches(1.0),
      fill=[BLUE, TEAL, NAVY, GREEN][i])
    T(sl, label, bx, Inches(4.25), Inches(2.4), Inches(0.45),
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, sub, bx, Inches(4.68), Inches(2.4), Inches(0.35),
      sz=10, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "→", bx + Inches(2.45), Inches(4.5), Inches(0.25), Inches(0.45),
          sz=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

T(sl, "Columbia University  ·  Digital Marketplaces  ·  Spring 2026",
  Inches(1), Inches(5.55), Inches(11.3), Inches(0.4),
  sz=13, color=RGBColor(0x7F, 0xA7, 0xCC), align=PP_ALIGN.CENTER)

notes(sl, "This project asks what happens when discovery in food delivery moves from platform apps to an LLM like ChatGPT. We build a simulation with 1,000 consumers, 50 restaurants, and 3 platforms, and test four payment mechanisms to ask: who benefits, who loses, and what design preserves consumer welfare?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Market Shift
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Discovery is moving upstream")

# Left: Traditional
R(sl, Inches(0.3), Inches(1.1), Inches(5.8), Inches(5.1), fill=SILVER,
  line=GREY, lw=Pt(1))
T(sl, "Traditional", Inches(0.3), Inches(1.12), Inches(5.8), Inches(0.48),
  sz=14, bold=True, color=GREY, align=PP_ALIGN.CENTER)
for i, (label, bg) in enumerate([
    ("Consumer", WHITE),
    ("Opens Platform App", WHITE),
    ("Sponsored platform feed", RGBColor(0xFF, 0xEC, 0xB3)),
    ("Order → fulfillment", WHITE),
]):
    R(sl, Inches(0.5), Inches(1.72) + i * Inches(1.08), Inches(5.4), Inches(0.82), fill=bg)
    T(sl, label, Inches(0.55), Inches(1.77) + i * Inches(1.08), Inches(5.3), Inches(0.72),
      sz=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "↓", Inches(2.9), Inches(2.54) + i * Inches(1.08), Inches(0.4), Inches(0.3),
          sz=16, color=GREY, align=PP_ALIGN.CENTER)

# VS
T(sl, "VS", Inches(6.17), Inches(3.3), Inches(1.0), Inches(0.65),
  sz=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Right: LLM-mediated
R(sl, Inches(7.2), Inches(1.1), Inches(5.8), Inches(5.1), fill=WHITE,
  line=BLUE, lw=Pt(2.5))
T(sl, "LLM-Mediated  (emerging)", Inches(7.2), Inches(1.12), Inches(5.8), Inches(0.48),
  sz=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
for i, (label, bg) in enumerate([
    ("Consumer", WHITE),
    ("LLM aggregates all platforms", RGBColor(0xBB, 0xDE, 0xFB)),
    ("Top-5 shortlist  (who pays?)", RGBColor(0xFF, 0xEC, 0xB3)),
    ("Platform fulfillment → Order", WHITE),
]):
    R(sl, Inches(7.4), Inches(1.72) + i * Inches(1.08), Inches(5.4), Inches(0.82), fill=bg)
    T(sl, label, Inches(7.45), Inches(1.77) + i * Inches(1.08), Inches(5.3), Inches(0.72),
      sz=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "↓", Inches(10.0), Inches(2.54) + i * Inches(1.08), Inches(0.4), Inches(0.3),
          sz=16, color=BLUE, align=PP_ALIGN.CENTER)

R(sl, Inches(0.3), Inches(6.3), Inches(12.73), Inches(0.5), fill=NAVY)
T(sl, "The platform still fulfills the order and earns commission + delivery fee.  The LLM only changes discovery.",
  Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.42),
  sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(sl, "Traditional delivery platforms control the ranking surface — what you see is what they choose to show you, with sponsored placement shaping the feed. In an LLM-mediated world (e.g., ChatGPT + Uber Eats), the AI becomes the discovery layer. Critically, platforms remain the fulfillment providers: they still process payments, send delivery drivers, and collect commissions. The LLM only changes who decides which options the consumer sees first. Transition: once the LLM controls discovery, the question is what it should be paid to optimize.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research Question
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "What should the LLM be paid to optimize?")

# Three tension boxes
for i, (label, bullets, color) in enumerate([
    ("Consumer Fit",
     "better matching\nlower prices\nfaster delivery", TEAL),
    ("LLM Monetization",
     "clicks · orders\nfulfilled intent\nplatform fees", RGBColor(0xB8, 0x86, 0x00)),
    ("Platform Revenue",
     "commission share\norder concentration\nWTP for LLM", NAVY),
]):
    bx = Inches(0.3) + i * Inches(4.35)
    R(sl, bx, Inches(1.15), Inches(4.1), Inches(2.6), fill=color)
    T(sl, label, bx, Inches(1.2), Inches(4.1), Inches(0.65),
      sz=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx + Inches(0.2), Inches(1.88), Inches(3.7), Pt(1.5), fill=GOLD)
    T(sl, bullets, bx, Inches(2.0), Inches(4.1), Inches(1.5),
      sz=13, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)

# Connecting arrows
for ax in [Inches(4.46), Inches(8.81)]:
    T(sl, "↔", ax, Inches(2.1), Inches(0.2), Inches(0.5),
      sz=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

# Main question block
R(sl, Inches(0.3), Inches(4.0), Inches(12.73), Inches(2.3), fill=NAVY)
T(sl, "Main question:", Inches(0.5), Inches(4.1), Inches(3.5), Inches(0.42),
  sz=14, italic=True, color=LIGHT)
T(sl, "How do different LLM payment mechanisms affect consumer welfare, platform profits, and market concentration?",
  Inches(0.5), Inches(4.55), Inches(12.3), Inches(1.5),
  sz=22, bold=True, color=WHITE)

notes(sl, "The three-way tension is the core of the project. Neutral LLM ranking maximizes consumer fit but gives platforms no reason to pay. Paid ranking can restore platform incentives but may distort recommendations. The mechanism design question is: which payment structure aligns LLM incentives with consumer welfare while still making participation viable for platforms? Transition: to answer this we build a simulation of the full market.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Actors and Choice Sets
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Four actors, two discovery worlds")

# Column headers
_cw = [Inches(2.0), Inches(5.2), Inches(5.53)]
th(sl, ["Actor", "What They Have", "What They Choose / Affect"],
   Inches(0.3), Inches(1.15), _cw)

# Actor rows
_rows = [
    ("Consumers", TEAL,
     "cuisine needs · budget · max delivery time\nplatform loyalty · LLM trust · sensitivities",
     "buy / not buy  ·  click or ignore LLM shortlist"),
    ("Restaurants", GREEN,
     "cuisine · quality (2–5 ★) · menu price\nprep time · cost ratio · 1–3 platform listings",
     "appear as offers on platforms"),
    ("Platforms", BLUE,
     "commission rate · delivery fee range\nlogistics time · internal ranking rule",
     "fulfill orders · sponsor LLM in paid regimes"),
    ("LLM", PURPLE,
     "ranking formula · payment mechanism\naggregates all platform offers",
     "returns top-5 shortlist to consumer"),
]
for i, (actor, color, have, choose) in enumerate(_rows):
    by = Inches(1.55) + i * Inches(1.2)
    bg = RGBColor(0xF4, 0xF7, 0xFA) if i % 2 == 0 else WHITE
    R(sl, Inches(0.3), by, Inches(2.0), Inches(1.1), fill=color)
    T(sl, actor, Inches(0.3), by + Inches(0.28), Inches(2.0), Inches(0.55),
      sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, Inches(2.3), by, Inches(5.2), Inches(1.1), fill=bg, line=GREY, lw=Pt(0.5))
    T(sl, have, Inches(2.38), by + Inches(0.12), Inches(5.04), Inches(0.88),
      sz=11, color=DARK)
    R(sl, Inches(7.5), by, Inches(5.53), Inches(1.1), fill=bg, line=GREY, lw=Pt(0.5))
    T(sl, choose, Inches(7.58), by + Inches(0.22), Inches(5.37), Inches(0.65),
      sz=11, color=DARK)

T(sl, "In the No-LLM world, consumers search their 2 most-loyal platforms.  In the LLM world, the AI aggregates all platforms and returns a ranked top-5 shortlist.",
  Inches(0.3), Inches(6.45), Inches(12.73), Inches(0.52),
  sz=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)

notes(sl, "Four actors: 1,000 consumers with heterogeneous preferences; 50 restaurants across 9 cuisines; three platforms (QuickEats 20% commission, FoodRush 25%, NomNom 15%); one LLM. The key model distinction is the choice set: no-LLM consumers see only their loyal platforms' feeds; LLM consumers get a cross-platform shortlist. Transition: let's see how the simulation puts these actors together.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Simulation Pipeline
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "We simulate a marketplace, then let consumers search")

_pipeline = [
    ("Generate\nmarket",     "3 platforms\n50 restaurants\nall offers",     BLUE),
    ("Generate\nconsumers",  "1,000 consumers\npreferences\nloyalty · trust", TEAL),
    ("Sample\nintent",       "cuisine · budget\nmax time\nper session",      GREEN),
    ("Build\nchoice set",    "No-LLM: 2 platforms\nLLM: top-5\ncross-platform", RGBColor(0xB8, 0x86, 0x00)),
    ("Simulate\nchoice",     "MNL logit\nclick → buy\nor abandon",           ORANGE),
    ("Record &\naggregate",  "utility · IF\nplatform net\nLLM revenue",      RGBColor(0x6C, 0x35, 0x8A)),
]
BOX_W = Inches(1.95)
STEP  = Inches(2.12)
for i, (name, detail, color) in enumerate(_pipeline):
    bx = Inches(0.3) + i * STEP
    R(sl, bx, Inches(1.15), BOX_W, Inches(4.8), fill=color)
    T(sl, name, bx, Inches(1.2), BOX_W, Inches(0.75),
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx + Inches(0.1), Inches(1.98), BOX_W - Inches(0.2), Pt(1.5), fill=GOLD)
    T(sl, detail, bx, Inches(2.1), BOX_W, Inches(2.8),
      sz=11, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    if i < 5:
        T(sl, "→", bx + BOX_W + Inches(0.04), Inches(2.8), Inches(0.14), Inches(0.5),
          sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Loop note
R(sl, Inches(0.3), Inches(6.1), Inches(12.73), Inches(0.62), fill=INK)
T(sl, "Repeat × 10 Monte Carlo runs  ·  10,000 consumer sessions per scenario  ·  Metrics averaged across runs",
  Inches(0.5), Inches(6.17), Inches(12.3), Inches(0.48),
  sz=13, color=LIGHT, align=PP_ALIGN.CENTER)

notes(sl, "The simulation builds a fresh market in each Monte Carlo run: platforms, restaurants, and 1,000 consumers are generated from their respective distributions. Within each run, every consumer draws a session-level intent, gets a choice set, makes a purchase decision, and records outcomes. We average across 10 runs to reduce noise from random initialization. Transition: the key structural choice is how the choice set is built — which leads to the two-worlds comparison.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Two Search Worlds
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "The experiment changes who controls discovery")

# No-LLM panel
R(sl, Inches(0.3), Inches(1.1), Inches(5.8), Inches(5.2), fill=WHITE,
  line=GREY, lw=Pt(1.5))
T(sl, "No-LLM  (baseline)", Inches(0.3), Inches(1.12), Inches(5.8), Inches(0.48),
  sz=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
for i, (label, bg) in enumerate([
    ("Consumer states meal intent", WHITE),
    ("Opens 2 most-loyal platforms\n(cuisine filter applied)", WHITE),
    ("Sees platform-ranked feed\n(quality · commission · promo · sponsored)", RGBColor(0xFF, 0xEC, 0xB3)),
    ("Purchase or no purchase", SILVER),
]):
    R(sl, Inches(0.5), Inches(1.72) + i * Inches(1.1), Inches(5.4), Inches(0.88), fill=bg)
    T(sl, label, Inches(0.55), Inches(1.77) + i * Inches(1.1), Inches(5.3), Inches(0.8),
      sz=12, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "↓", Inches(2.9), Inches(2.6) + i * Inches(1.1), Inches(0.4), Inches(0.3),
          sz=16, color=GREY, align=PP_ALIGN.CENTER)

# Callout
R(sl, Inches(0.4), Inches(6.0), Inches(5.6), Inches(0.42), fill=GREY)
T(sl, "Fragmented search — 2 platforms only",
  Inches(0.4), Inches(6.03), Inches(5.6), Inches(0.38),
  sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# VS
T(sl, "VS", Inches(6.2), Inches(3.25), Inches(0.95), Inches(0.65),
  sz=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# LLM panel
R(sl, Inches(7.2), Inches(1.1), Inches(5.8), Inches(5.2), fill=WHITE,
  line=BLUE, lw=Pt(2.5))
T(sl, "LLM  (treatment)", Inches(7.2), Inches(1.12), Inches(5.8), Inches(0.48),
  sz=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
for i, (label, bg) in enumerate([
    ("Consumer states meal intent", WHITE),
    ("LLM aggregates all 3 platforms\n(sees every offer)", RGBColor(0xBB, 0xDE, 0xFB)),
    ("Top-5 shortlist returned\n(intent fit ± sponsorship bias)", RGBColor(0xFF, 0xEC, 0xB3)),
    ("Click → purchase or abandon", SILVER),
]):
    R(sl, Inches(7.4), Inches(1.72) + i * Inches(1.1), Inches(5.4), Inches(0.88), fill=bg)
    T(sl, label, Inches(7.45), Inches(1.77) + i * Inches(1.1), Inches(5.3), Inches(0.8),
      sz=12, color=DARK, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "↓", Inches(10.0), Inches(2.6) + i * Inches(1.1), Inches(0.4), Inches(0.3),
          sz=16, color=BLUE, align=PP_ALIGN.CENTER)

R(sl, Inches(7.3), Inches(6.0), Inches(5.6), Inches(0.42), fill=BLUE)
T(sl, "Cross-platform aggregation — all 3 platforms",
  Inches(7.3), Inches(6.03), Inches(5.6), Inches(0.38),
  sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(sl, "The no-LLM baseline is a strong one: consumers are loyal to specific platforms and use cuisine filters. They are not random choosers. The LLM advantage comes purely from cross-platform aggregation — seeing all 50 restaurants across all 3 platforms rather than just those on 2 apps. The payment regime determines whether the shortlist is purely intent-ranked or partially biased toward sponsors. Transition: before showing results, we define how we measure outcomes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — How An Order Is Scored
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Each order is scored from four perspectives")

_metrics = [
    ("Consumer", [
        "Intent fulfilment  ← primary",
        "Conversion rate",
        "Avg price paid",
        "Avg delivery time",
        "Utility",
    ], NAVY, Inches(0.3), Inches(1.15)),
    ("Platform", [
        "Net revenue",
        "Order share by platform",
        "Net vs no-LLM baseline",
        "WTP for LLM",
    ], BLUE, Inches(6.98), Inches(1.15)),
    ("Restaurant", [
        "Profit per order",
        "Exposure concentration",
        "(HHI · Gini)",
    ], TEAL, Inches(0.3), Inches(4.0)),
    ("LLM", [
        "Total revenue",
        "Click-through rate",
        "Shortlist relevance",
        "Revenue per order",
    ], PURPLE, Inches(6.98), Inches(4.0)),
]
for name, items, color, bx, by in _metrics:
    R(sl, bx, by, Inches(6.35), Inches(2.5), fill=color)
    T(sl, name, bx, by + Inches(0.07), Inches(6.35), Inches(0.52),
      sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx, by + Inches(0.59), Inches(6.35), Pt(1.5), fill=GOLD)
    for j, item in enumerate(items):
        T(sl, "·  " + item, bx + Inches(0.15), by + Inches(0.7) + j * Inches(0.38),
          Inches(6.1), Inches(0.36), sz=12, color=WHITE)

# Formula strip
R(sl, Inches(0.3), Inches(6.6), Inches(12.73), Inches(0.55), fill=INK)
T(sl, "Platform net  =  commission  +  delivery fee  −  promo subsidy  −  LLM payment",
  Inches(0.5), Inches(6.66), Inches(12.3), Inches(0.42),
  sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

notes(sl, "We track outcomes for all four sides of the market. The primary consumer welfare metric is intent fulfilment: did the completed order match the consumer's stated cuisine, budget, delivery time, and quality preferences? This is cleaner than utility because it measures whether the stated request was satisfied, not just whether the consumer was happy. Platform net revenue is the key supply-side metric — including the LLM payment cost so we can directly compare profitability across regimes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Intent Fulfilment
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "Consumer welfare: did the user get what they asked for?")

T(sl, "Intent Fulfilment  =", Inches(0.4), Inches(1.1), Inches(4.0), Inches(0.65),
  sz=22, bold=True, color=GOLD)

_factors = [
    ("Cuisine\nmatch", "Wrong cuisine\n→ score = 0\n\nExact match\n→ score = 1.0", TEAL),
    ("×\nBudget\nfit", "Over budget\n→ graded penalty\n\nWithin budget\n→ score = 1.0", BLUE),
    ("×\nTime\nfit", "Over time limit\n→ penalty/min late\n\nOn time\n→ score = 1.0", GREEN),
    ("×\nQuality\nadjustment", "2.0 stars → 0.60\n3.0 stars → 0.73\n4.5 stars → 0.93\n5.0 stars → 1.00", ORANGE),
]
for i, (name, detail, color) in enumerate(_factors):
    bx = Inches(0.3) + i * Inches(3.28)
    R(sl, bx, Inches(1.85), Inches(3.0), Inches(4.3), fill=color)
    T(sl, name, bx, Inches(1.9), Inches(3.0), Inches(1.0),
      sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx + Inches(0.1), Inches(2.9), Inches(2.8), Pt(1.5), fill=GOLD)
    T(sl, detail, bx + Inches(0.1), Inches(3.0), Inches(2.8), Inches(2.8),
      sz=12, color=WHITE, align=PP_ALIGN.CENTER)

T(sl, "Multiplicative — a zero on any factor (wrong cuisine, or extreme price/time overrun) collapses fulfilment to zero.",
  Inches(0.4), Inches(6.35), Inches(12.73), Inches(0.48),
  sz=13, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

notes(sl, "Intent fulfilment is multiplicative: all four factors must be satisfied for a high score. A wrong cuisine immediately collapses the score to zero regardless of price or quality. Price and time violations are graded penalties — slightly over budget or slightly late still earns partial credit. Quality adjusts the maximum possible score but never dominates the other factors. This metric is applied only to completed orders, not to sessions that end without a purchase. Transition: to make this concrete, let's walk through one consumer's session across all four regimes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Concrete Walkthrough
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "One consumer shows the mechanism")

# Consumer intent card
R(sl, Inches(0.3), Inches(1.1), Inches(2.9), Inches(5.8), fill=NAVY)
T(sl, "CONSUMER\nINTENT", Inches(0.3), Inches(1.2), Inches(2.9), Inches(0.85),
  sz=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
R(sl, Inches(0.4), Inches(2.05), Inches(2.7), Pt(1.5), fill=GOLD)
for j, (label, val) in enumerate([
    ("Cuisine", "Korean"),
    ("Budget", "$50"),
    ("Max time", "45 min"),
    ("Loyal to", "FoodRush\n+ NomNom"),
]):
    T(sl, label, Inches(0.45), Inches(2.22) + j * Inches(0.98), Inches(2.6), Inches(0.32),
      sz=10, italic=True, color=LIGHT)
    T(sl, val, Inches(0.45), Inches(2.52) + j * Inches(0.98), Inches(2.6), Inches(0.55),
      sz=18, bold=True, color=WHITE)
    if j < 3:
        R(sl, Inches(0.5), Inches(3.08) + j * Inches(0.98), Inches(2.4), Pt(1),
          fill=RGBColor(0x2A, 0x3E, 0x5C))

# Comparison table
_tw = [Inches(1.65), Inches(3.1), Inches(2.9), Inches(2.53)]
_tx = Inches(3.4)
th(sl, ["World", "What Happens", "Chosen Offer", "Intent Fulfil."],
   _tx, Inches(1.1), _tw)

_rows9 = [
    ("No-LLM",     "Checks FoodRush +\nNomNom only",          "NomNom  ·  51 min", "0.763", GREY,   WHITE),
    ("LLM Neutral","Sees all 3 platforms",                     "NomNom  ·  29 min", "0.917", TEAL,   WHITE),
    ("LLM CPC",    "FoodRush gets\nhighest paid boost",        "FoodRush  ·  42 min","0.866", ORANGE, WHITE),
    ("LLM CPFI",   "Boost scaled by\nintent fit",              "QuickEats  ·  20 min","0.931", GREEN,  WHITE),
]
for i, (world, what, offer, IF, hdr_color, _) in enumerate(_rows9):
    by = Inches(1.5) + i * Inches(1.25)
    bg = RGBColor(0xF4, 0xF7, 0xFA) if i % 2 == 0 else WHITE
    R(sl, _tx, by, _tw[0], Inches(1.15), fill=hdr_color)
    T(sl, world, _tx + Inches(0.06), by + Inches(0.2), _tw[0] - Inches(0.12), Inches(0.75),
      sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, _tx + _tw[0], by, _tw[1], Inches(1.15), fill=bg, line=GREY, lw=Pt(0.5))
    T(sl, what, _tx + _tw[0] + Inches(0.08), by + Inches(0.15),
      _tw[1] - Inches(0.16), Inches(0.88), sz=11, color=DARK)
    R(sl, _tx + _tw[0] + _tw[1], by, _tw[2], Inches(1.15), fill=bg, line=GREY, lw=Pt(0.5))
    T(sl, offer, _tx + _tw[0] + _tw[1] + Inches(0.08), by + Inches(0.28),
      _tw[2] - Inches(0.16), Inches(0.6), sz=11, color=DARK)
    # IF score
    if_color = GREEN if float(IF) >= 0.91 else (TEAL if float(IF) >= 0.85 else (ORANGE if float(IF) >= 0.80 else GREY))
    R(sl, _tx + _tw[0] + _tw[1] + _tw[2], by, _tw[3], Inches(1.15), fill=if_color)
    T(sl, IF, _tx + _tw[0] + _tw[1] + _tw[2], by + Inches(0.25), _tw[3], Inches(0.65),
      sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(sl, "This is a stylized walkthrough illustrating the model mechanics. In the no-LLM world, the consumer opens FoodRush and NomNom — their two most-loyal platforms. The best Korean offer available through those two is on NomNom at 51 minutes — above the 45-minute limit, so intent fulfilment is penalized for time. The neutral LLM aggregates all three platforms and finds a NomNom option at only 29 minutes — well within the time limit — raising intent fulfilment to 0.917. Under CPC, FoodRush gets the highest paid boost because it pays the most per click, routing to a FoodRush option at 42 minutes with moderate fulfilment. Under CPFI, the boost is multiplied by the predicted intent fulfilment score, so it only amplifies genuinely good matches — QuickEats at 20 minutes scores highest at 0.931. The key insight: CPC rewards the highest bidder; CPFI rewards the best match.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Experiment Roadmap
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "Four experiments build one argument")

for i, (num, title, question, color) in enumerate([
    ("1", "No-LLM\nvs LLM Neutral",
     "Does AI improve\nconsumer matching?", TEAL),
    ("2", "Neutral vs CPC\nvs CPA vs CPFI",
     "How should\nthe LLM be paid?", BLUE),
    ("3", "Sponsorship bias\nsweep  0.0 → 1.0",
     "How much paid bias\nis too much?", ORANGE),
    ("4", "Consumer trust\nsweep",
     "Does adoption change\nquality or volume?", GREEN),
]):
    bx = Inches(0.35) + i * Inches(3.25)
    R(sl, bx, Inches(1.1), Inches(3.05), Inches(5.7), fill=color)
    T(sl, num, bx, Inches(1.15), Inches(3.05), Inches(1.2),
      sz=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, title, bx, Inches(2.4), Inches(3.05), Inches(1.35),
      sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx + Inches(0.15), Inches(3.8), Inches(2.75), Pt(1.5), fill=GOLD)
    T(sl, question, bx + Inches(0.1), Inches(3.95), Inches(2.9), Inches(1.55),
      sz=14, italic=True, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)

notes(sl, "The experiments are sequential and build one argument. Experiment 1 establishes the consumer-platform tension: neutral LLM helps consumers but hurts platforms. Experiment 2 tests four payment mechanisms to resolve that tension. Experiment 3 finds the threshold where paid bias becomes harmful. Experiment 4 tests whether the results depend on how much consumers trust the LLM.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Experiment 1: Consumer Results
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "Exp 1: Neutral LLM improves consumer matching")

cd = ChartData()
cd.categories = ["No-LLM\n(direct search)", "LLM Neutral"]
cd.add_series("Intent Fulfilment", (0.683, 0.747))
chart = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                  Inches(0.8), Inches(1.0), Inches(6.5), Inches(5.5))
chart.has_legend = False
set_bar_colors(chart, [[RGBColor(0x5D, 0x8A, 0xB0), GREEN]])
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.font.size = Pt(18)
chart.plots[0].data_labels.font.bold = True
chart.plots[0].data_labels.font.color.rgb = WHITE
chart.value_axis.has_major_gridlines = False
chart.value_axis.minimum_scale = 0.60
chart.value_axis.maximum_scale = 0.80
chart.value_axis.tick_labels.font.size = Pt(11)
chart.value_axis.tick_labels.font.color.rgb = LIGHT
chart.category_axis.tick_labels.font.size = Pt(14)
chart.category_axis.tick_labels.font.color.rgb = WHITE

for i, (num, lbl, bg) in enumerate([
    ("−$1.71",   "avg price paid",   TEAL),
    ("−1.7 min", "delivery time",    BLUE),
    ("+1.2 pp",  "conversion rate",  RGBColor(0x5D, 0x6D, 0x7E)),
]):
    bx, by = Inches(7.6), Inches(1.0) + i * Inches(1.85)
    R(sl, bx, by, Inches(5.4), Inches(1.65), fill=bg)
    T(sl, num, bx, by + Inches(0.1), Inches(5.4), Inches(0.85),
      sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, lbl, bx, by + Inches(1.05), Inches(5.4), Inches(0.45),
      sz=15, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)

notes(sl, "The neutral LLM — no paid bias — raises intent fulfilment from 0.683 to 0.747, a 9.3% improvement, even after giving the no-LLM baseline both platform loyalty and cuisine filters. The gain comes entirely from cross-platform aggregation: consumers find better-fit offers that their loyal platforms didn't carry. Average price falls $1.71, delivery time falls 1.7 minutes, conversion rises 1.2pp. Transition: but there is a problem on the platform side.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Experiment 1: Platform Results
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Exp 1: But neutral LLM weakens platform incentives")

R(sl, Inches(0.4), Inches(1.1), Inches(5.6), Inches(5.6), fill=SILVER)
T(sl, "No-LLM",  Inches(0.4), Inches(1.15), Inches(5.6), Inches(0.5),
  sz=15, bold=True, color=GREY, align=PP_ALIGN.CENTER)
T(sl, "$69,366", Inches(0.4), Inches(2.1),  Inches(5.6), Inches(1.8),
  sz=68, bold=True, color=DARK, align=PP_ALIGN.CENTER)
T(sl, "aggregate platform net revenue", Inches(0.4), Inches(3.95), Inches(5.6), Inches(0.4),
  sz=13, color=GREY, align=PP_ALIGN.CENTER)

T(sl, "→", Inches(6.15), Inches(3.5), Inches(1.0), Inches(0.7),
  sz=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

R(sl, Inches(7.3), Inches(1.1), Inches(5.6), Inches(5.6), fill=RGBColor(0xFD, 0xED, 0xEB))
T(sl, "LLM Neutral", Inches(7.3), Inches(1.15), Inches(5.6), Inches(0.5),
  sz=15, bold=True, color=RED, align=PP_ALIGN.CENTER)
T(sl, "$67,884", Inches(7.3), Inches(2.1), Inches(5.6), Inches(1.8),
  sz=68, bold=True, color=RED, align=PP_ALIGN.CENTER)
T(sl, "aggregate platform net revenue", Inches(7.3), Inches(3.95), Inches(5.6), Inches(0.4),
  sz=13, color=RED, align=PP_ALIGN.CENTER)
R(sl, Inches(7.3), Inches(4.6), Inches(5.6), Inches(0.55), fill=RED)
T(sl, "−$1,483  ·  Negative platform WTP", Inches(7.3), Inches(4.62), Inches(5.6), Inches(0.5),
  sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "Better matches route away from high-commission options",
  Inches(7.3), Inches(5.25), Inches(5.6), Inches(0.45),
  sz=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)

notes(sl, "The neutral LLM is good for consumers but bad for platforms. When the LLM ranks purely by consumer intent, it routes orders toward best-fit restaurants — which are often lower-commission or less sponsored than what a platform's own feed would surface. Aggregate platform net falls $1,483. Restaurant profit per order also falls 8.6%. This means platforms have no financial incentive to pay for neutral LLM access. Paid regimes exist to fix this — they change routing in ways that can make platforms willing to participate. Transition: but how exactly does that work?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Why Paid Regimes Can Restore Incentives
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Payment is a cost — but routing is a benefit")

# Equation strip
R(sl, Inches(0.3), Inches(1.1), Inches(12.73), Inches(0.95), fill=NAVY)
T(sl, "Extra commission from redirected orders   −   LLM payment   =   platform net gain",
  Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.75),
  sz=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Mechanism flow
for i, (label, detail, color) in enumerate([
    ("Platform\nsponsors LLM",   "pays per click\nor per order", NAVY),
    ("LLM ranking\nshifts",      "sponsor's offers\nrise in shortlist", BLUE),
    ("More orders\nto sponsor",  "from consumers\nwho would not\nhave seen them", TEAL),
    ("Extra commission\n> LLM fee", "net gain for\nplatform", GREEN),
]):
    bx = Inches(0.3) + i * Inches(3.25)
    R(sl, bx, Inches(2.2), Inches(3.0), Inches(2.3), fill=color)
    T(sl, label, bx, Inches(2.25), Inches(3.0), Inches(0.85),
      sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, detail, bx, Inches(3.1), Inches(3.0), Inches(1.1),
      sz=11, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl, "→", bx + Inches(3.07), Inches(3.1), Inches(0.16), Inches(0.45),
          sz=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Nuance callout
R(sl, Inches(0.3), Inches(4.7), Inches(6.35), Inches(1.65), fill=RGBColor(0xFD, 0xED, 0xEB))
T(sl, "Important nuance", Inches(0.45), Inches(4.8), Inches(6.05), Inches(0.4),
  sz=13, bold=True, color=RED)
T(sl, "This is a competitive sponsored-discovery game, not a cooperative agreement. Sponsors may gain; low-bid rivals (e.g., NomNom) can lose traffic as a competitive externality.",
  Inches(0.45), Inches(5.22), Inches(6.05), Inches(1.0), sz=12, color=DARK, wrap=True)

R(sl, Inches(6.98), Inches(4.7), Inches(6.05), Inches(1.65), fill=INK)
T(sl, "WTP calibration:", Inches(7.1), Inches(4.8), Inches(5.8), Inches(0.38),
  sz=12, bold=True, color=GOLD)
T(sl, "LLM captures 50% of positive gross sponsor-side surplus. Rates are calibrated to value created by routing — not arbitrary flat fees.",
  Inches(7.1), Inches(5.2), Inches(5.75), Inches(1.0), sz=11, color=LIGHT, wrap=True)

notes(sl, "Platforms pay the LLM voluntarily because routing bias brings more orders. The extra commission revenue can exceed the LLM payment. We calibrate by first measuring gross platform surplus when routing is on but payments are zero, then charging the LLM 50% of that surplus. This ensures payments are tied to value created. Critically, this is not a cooperative equilibrium: high-bid platforms gain at the expense of low-bid rivals, similar to search advertising. Transition: this creates the mechanism design question — which payment trigger best aligns the LLM's incentive with consumer welfare?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Payment Mechanisms
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "The payment trigger determines the incentive")

# Table header
_pw = [Inches(2.0), Inches(2.0), Inches(5.53), Inches(3.0)]
th(sl, ["Regime", "LLM Paid When...", "Incentive Created", "Consumer Risk"],
   Inches(0.3), Inches(1.15), _pw)

_prows = [
    ("Neutral", "Never",
     "Rank purely by consumer intent fit",
     "None — but platforms have negative WTP",
     TEAL),
    ("CPC", "User clicks",
     "Maximize attention to sponsored offers",
     "Recommends high-bidders even if poor fit",
     RED),
    ("CPA", "User orders",
     "Maximize conversion from sponsored offers",
     "Moderate — off-intent orders less likely",
     ORANGE),
    ("CPFI", "Order × intent score",
     "Maximize fulfilled consumer intent",
     "Low — boost only helps when offer fits",
     GREEN),
]
for i, (regime, when, incentive, risk, color) in enumerate(_prows):
    by = Inches(1.55) + i * Inches(1.18)
    R(sl, Inches(0.3), by, Inches(2.0), Inches(1.08), fill=color)
    T(sl, regime, Inches(0.3), by + Inches(0.25), Inches(2.0), Inches(0.58),
      sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (text, xoff, ww) in enumerate([
        (when,      _pw[0],            _pw[1]),
        (incentive, _pw[0] + _pw[1],   _pw[2]),
        (risk,      sum(_pw[:3]),       _pw[3]),
    ]):
        bg = RGBColor(0xF4, 0xF7, 0xFA) if i % 2 == 0 else WHITE
        R(sl, Inches(0.3) + xoff, by, ww, Inches(1.08), fill=bg, line=GREY, lw=Pt(0.5))
        T(sl, text, Inches(0.3) + xoff + Inches(0.08), by + Inches(0.15),
          ww - Inches(0.16), Inches(0.8), sz=11, color=DARK, wrap=True)

R(sl, Inches(0.3), Inches(6.4), Inches(12.73), Inches(0.48), fill=RGBColor(0x10, 0x1E, 0x33))
T(sl, "CPFI = Cost Per Fulfilled Intent  ·  payment scales with how well the order satisfied the consumer's stated goal",
  Inches(0.5), Inches(6.46), Inches(12.3), Inches(0.38),
  sz=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

notes(sl, "This is the mechanism-design heart of the project. CPC rewards attention — the LLM is paid for clicks regardless of whether the consumer is satisfied. CPA rewards completed orders — better than CPC but still doesn't reward satisfaction. CPFI rewards fulfilled intent — the LLM's payment is proportional to how well the order matched the consumer's stated cuisine, budget, and time requirements. Under CPFI, off-intent sponsored offers get little amplification even if the platform bids high, because the payment would be small. Transition: let's see how these mechanisms affect actual outcomes.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Experiment 2: Consumer Welfare
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "Exp 2: CPC distorts; CPFI preserves fit")

cd = ChartData()
cd.categories = ["Neutral", "CPC", "CPA", "CPFI"]
cd.add_series("Intent Fulfilment", (0.747, 0.716, 0.726, 0.747))
chart = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                  Inches(0.5), Inches(0.9), Inches(7.5), Inches(5.9))
chart.has_legend = False
set_bar_colors(chart, [[TEAL, RED, ORANGE, GREEN]])
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.font.size = Pt(20)
chart.plots[0].data_labels.font.bold = True
chart.plots[0].data_labels.font.color.rgb = WHITE
chart.value_axis.has_major_gridlines = False
chart.value_axis.minimum_scale = 0.68
chart.value_axis.maximum_scale = 0.77
chart.value_axis.tick_labels.font.color.rgb = LIGHT
chart.value_axis.tick_labels.font.size = Pt(11)
chart.category_axis.tick_labels.font.size = Pt(16)
chart.category_axis.tick_labels.font.bold = True
chart.category_axis.tick_labels.font.color.rgb = WHITE

for i, (regime, desc, sr, color) in enumerate([
    ("Neutral", "No payment — baseline",         "Shortlist rel.  0.973", TEAL),
    ("CPC",     "Per click — highest distortion", "Shortlist rel.  0.929", RED),
    ("CPA",     "Per order — moderate",           "Shortlist rel.  0.946", ORANGE),
    ("CPFI",    "Per fulfilled intent — best",    "Shortlist rel.  0.981", GREEN),
]):
    R(sl, Inches(8.2), Inches(1.2) + i * Inches(1.2), Inches(4.8), Inches(1.0), fill=color)
    T(sl, regime, Inches(8.35), Inches(1.25) + i * Inches(1.2), Inches(4.5), Inches(0.38),
      sz=17, bold=True, color=WHITE)
    T(sl, desc + "  ·  " + sr,
      Inches(8.35), Inches(1.6) + i * Inches(1.2), Inches(4.5), Inches(0.38),
      sz=10, italic=True, color=RGBColor(0xDD, 0xDD, 0xDD))

notes(sl, "CPC causes the largest distortion because it rewards clicks regardless of satisfaction. FoodRush pays $0.80 per click vs NomNom's $0.30, so the LLM biases toward FoodRush even when its restaurants don't fit the consumer's intent — hence the shortlist relevance drop from 0.973 to 0.929. CPFI matches neutral intent fulfilment (0.747) because each offer's sponsorship bonus is multiplied by its predicted intent fulfilment score — off-intent sponsored options get little amplification. CPFI shortlist relevance is actually above neutral (0.981) because the bias now slightly rewards well-fitting offers.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Experiment 2: Platform Revenue
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Exp 2: Paid regimes restore aggregate platform profitability")

cd = ChartData()
cd.categories = ["No-LLM", "LLM\nNeutral", "LLM\nCPC", "LLM\nCPA", "LLM\nCPFI"]
cd.add_series("Platform Net Revenue ($)", (69366, 67884, 72468, 72239, 71309))
chart = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                  Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.6))
chart.has_legend = False
set_bar_colors(chart, [[GREY, RED, NAVY, BLUE, GREEN]])
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.font.size = Pt(13)
chart.plots[0].data_labels.font.bold = True
chart.plots[0].data_labels.font.color.rgb = DARK
chart.value_axis.has_major_gridlines = False
chart.value_axis.minimum_scale = 65000
chart.value_axis.maximum_scale = 75000
chart.value_axis.number_format = '"$"#,##0'
chart.value_axis.tick_labels.font.size = Pt(11)
chart.category_axis.tick_labels.font.size = Pt(13)
chart.category_axis.tick_labels.font.bold = True

R(sl, Inches(9.1), Inches(1.0), Inches(3.9), Inches(5.6), fill=INK)
T(sl, "vs no-LLM baseline", Inches(9.1), Inches(1.05), Inches(3.9), Inches(0.42),
  sz=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
for j, (name, delta, color) in enumerate([
    ("Neutral", "−$1,483", RED),
    ("CPC",     "+$3,101", GREEN),
    ("CPA",     "+$2,873", GREEN),
    ("CPFI",    "+$1,942", GREEN),
]):
    R(sl, Inches(9.2), Inches(1.55) + j * Inches(1.15), Inches(3.7), Inches(1.0), fill=color)
    T(sl, name + "   " + delta,
      Inches(9.2), Inches(1.58) + j * Inches(1.15), Inches(3.7), Inches(0.95),
      sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

T(sl, "Paid routing creates surplus even after LLM payments — because redirected orders generate commission revenue that exceeds the LLM fee.",
  Inches(0.4), Inches(6.72), Inches(12.5), Inches(0.52),
  sz=12, italic=True, color=GREY)

notes(sl, "All paid regimes exceed the no-LLM baseline because rates are calibrated to 50% of the value created by routing — platforms keep the other half. CPC generates the most ($72,468) because FoodRush's high commission rate makes redirected orders high-value. CPFI generates the least of the paid regimes ($71,309) because its alignment mechanism limits how much routing bias can be applied. But not every platform wins: NomNom loses traffic in all paid regimes as a competitive externality from FoodRush's sponsorship.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Experiment 2: Market Concentration
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Exp 2: CPC creates a winner-take-share outcome")

cd = ChartData()
cd.categories = ["Neutral", "CPC", "CPA", "CPFI"]
cd.add_series("FoodRush order share", (0.326, 0.667, 0.604, 0.465))
cd.add_series("NomNom order share",   (0.269, 0.038, 0.036, 0.128))
chart = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                  Inches(0.4), Inches(1.0), Inches(8.5), Inches(5.7))
chart.has_legend = True
chart.legend.font.size = Pt(13)
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = NAVY
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = TEAL
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.number_format = '0%'
chart.plots[0].data_labels.font.size = Pt(13)
chart.plots[0].data_labels.font.bold = True
chart.plots[0].data_labels.font.color.rgb = WHITE
chart.value_axis.has_major_gridlines = False
chart.value_axis.number_format = '0%'
chart.value_axis.tick_labels.font.size = Pt(11)
chart.category_axis.tick_labels.font.size = Pt(14)
chart.category_axis.tick_labels.font.bold = True

R(sl, Inches(9.1), Inches(1.0), Inches(3.9), Inches(2.6), fill=RED)
T(sl, "Under CPC", Inches(9.2), Inches(1.1), Inches(3.7), Inches(0.4),
  sz=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "66.7%", Inches(9.2), Inches(1.5), Inches(3.7), Inches(1.0),
  sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "FoodRush order share", Inches(9.2), Inches(2.5), Inches(3.7), Inches(0.4),
  sz=13, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

R(sl, Inches(9.1), Inches(3.8), Inches(3.9), Inches(2.9), fill=GREEN)
T(sl, "Under CPFI", Inches(9.2), Inches(3.9), Inches(3.7), Inches(0.4),
  sz=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "46.5%", Inches(9.2), Inches(4.3), Inches(3.7), Inches(1.0),
  sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "FoodRush order share", Inches(9.2), Inches(5.3), Inches(3.7), Inches(0.4),
  sz=13, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

notes(sl, "FoodRush pays $0.80 per click vs NomNom's $0.30 — 2.7x the rate. Under CPC, the unconditional bias bonus is large enough to dominate shortlist composition. FoodRush jumps from 33% to 67%; NomNom is nearly eliminated at 3.8%. This is a winner-take-share outcome driven entirely by payment design, not by FoodRush offering better restaurants. Under CPFI, concentration is dramatically lower: FoodRush holds 46.5%, NomNom retains 12.8%, because the bias only amplifies FoodRush when its restaurants actually fit the consumer's intent.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Experiment 3: Bias Sweep
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=NAVY)
footer(sl)
hdr(sl, "Exp 3: More paid bias eventually destroys the LLM advantage")

cd = ChartData()
cd.categories = ["0.0", "0.1", "0.2", "0.3\n(default)", "0.5", "0.7", "1.0"]
cd.add_series("Intent Fulfilment", (0.747, 0.741, 0.731, 0.716, 0.689, 0.663, 0.649))
chart = add_chart(sl, XL_CHART_TYPE.COLUMN_CLUSTERED, cd,
                  Inches(0.5), Inches(0.9), Inches(8.8), Inches(5.5))
chart.has_legend = False
set_bar_colors(chart, [[GREEN, GREEN, ORANGE, ORANGE, RED, RED, RED]])
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.font.size = Pt(13)
chart.plots[0].data_labels.font.bold = True
chart.plots[0].data_labels.font.color.rgb = WHITE
chart.value_axis.has_major_gridlines = False
chart.value_axis.minimum_scale = 0.62
chart.value_axis.maximum_scale = 0.77
chart.value_axis.tick_labels.font.color.rgb = LIGHT
chart.value_axis.tick_labels.font.size = Pt(11)
chart.category_axis.tick_labels.font.size = Pt(12)
chart.category_axis.tick_labels.font.color.rgb = WHITE

T(sl, "← Sponsorship bias strength →", Inches(1.5), Inches(6.52), Inches(6.0), Inches(0.35),
  sz=11, italic=True, color=LIGHT)

# No-LLM baseline annotation (approximate chart position for 0.683 on [0.62, 0.77] scale)
R(sl, Inches(0.6), Inches(4.25), Inches(8.5), Pt(2), fill=GOLD)
T(sl, "No-LLM baseline  =  0.683", Inches(0.65), Inches(3.95), Inches(4.5), Inches(0.3),
  sz=11, bold=True, color=GOLD)

R(sl, Inches(9.5), Inches(1.0), Inches(3.6), Inches(2.8), fill=RED)
T(sl, "Crossover", Inches(9.6), Inches(1.1), Inches(3.4), Inches(0.42),
  sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "At bias 1.0\nLLM score 0.649\n< no-LLM 0.683",
  Inches(9.6), Inches(1.55), Inches(3.4), Inches(1.95),
  sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

R(sl, Inches(9.5), Inches(4.0), Inches(3.6), Inches(2.7), fill=GREEN)
T(sl, "Recommended cap", Inches(9.6), Inches(4.1), Inches(3.4), Inches(0.38),
  sz=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "≤ 0.20", Inches(9.6), Inches(4.52), Inches(3.4), Inches(1.1),
  sz=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "preserves ~95%\nof welfare gains", Inches(9.6), Inches(5.62), Inches(3.4), Inches(0.65),
  sz=12, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

notes(sl, "Green bars are above the no-LLM baseline of 0.683. Orange bars are degraded but still better than no-LLM. Red bars fall below the no-LLM baseline — the LLM becomes worse than loyal direct platform search. The crossover happens between bias 0.7 and 1.0. Revenue gains are sharply diminishing above 0.5: raising bias from 0.5 to 1.0 adds only $358 in revenue while cutting intent fulfilment 4 percentage points. A cap at 0.10–0.20 preserves approximately 95% of consumer welfare gains at minimal revenue cost. This is analogous to disclosure-style regulation: cap the weight of paid placement, don't prohibit it. Transition: the final question is whether consumer adoption levels change these results.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Experiment 4: Trust Sweep
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "Exp 4: Trust drives volume, not recommendation quality")

cd = ChartData()
cd.categories = ["Very Low\n(0.0–0.2)", "Low\n(0.2–0.4)", "Medium\n(0.4–0.6)",
                  "High\n(0.6–0.8)", "Very High\n(0.8–1.0)"]
cd.add_series("LLM Conversion %",    (78.6, 80.2, 82.0, 83.7, 85.2))
cd.add_series("No-LLM Conversion %", (80.4, 80.4, 80.4, 80.4, 80.4))
chart = add_chart(sl, XL_CHART_TYPE.LINE, cd,
                  Inches(0.4), Inches(0.9), Inches(7.8), Inches(5.8))
chart.has_legend = True
chart.legend.font.size = Pt(12)
chart.series[0].format.line.color.rgb = BLUE
chart.series[0].format.line.width = Pt(3)
chart.series[1].format.line.color.rgb = GREY
chart.series[1].format.line.width = Pt(2)
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.show_value = True
chart.plots[0].data_labels.font.size = Pt(12)
chart.plots[0].data_labels.font.bold = True
chart.value_axis.has_major_gridlines = False
chart.value_axis.minimum_scale = 76
chart.value_axis.maximum_scale = 87
chart.value_axis.number_format = '0.0"%"'
chart.value_axis.tick_labels.font.size = Pt(11)
chart.category_axis.tick_labels.font.size = Pt(12)

R(sl, Inches(8.4), Inches(1.0), Inches(4.6), Inches(5.7), fill=NAVY)
T(sl, "Intent Fulfilment", Inches(8.5), Inches(1.1), Inches(4.4), Inches(0.45),
  sz=14, color=LIGHT, align=PP_ALIGN.CENTER)
T(sl, "FLAT", Inches(8.5), Inches(1.8), Inches(4.4), Inches(1.5),
  sz=64, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
T(sl, "0.747 → 0.749", Inches(8.5), Inches(3.35), Inches(4.4), Inches(0.5),
  sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, "across all trust levels", Inches(8.5), Inches(3.88), Inches(4.4), Inches(0.35),
  sz=12, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
R(sl, Inches(8.5), Inches(4.38), Inches(4.2), Inches(1.52), fill=RGBColor(0x2A, 0x3E, 0x5C))
T(sl, "Match quality = algorithm\nNot adoption rate",
  Inches(8.6), Inches(4.48), Inches(4.0), Inches(1.3),
  sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(sl, "Two effects emerge. Demand effect (strong): conversion rises from 78.6% to 85.2% as trust increases — consumers who trust the AI follow recommendations more often. Quality effect (negligible): intent fulfilment is essentially flat at 0.747–0.749 regardless of trust. The LLM's algorithm determines which options appear in the shortlist; trust only determines whether consumers act on them. Policy implication: even at low adoption, LLM recommendations are equally good matches. You don't need mass trust to benefit from good algorithmic design.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Final Takeaways
# ══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl)
hdr(sl, "The mechanism design choice is the whole game")

for i, (num, headline, body, bg) in enumerate([
    ("+9.3%",
     "Neutral LLM improves matching",
     "but platforms lose aggregate net revenue — negative WTP for neutral access",
     NAVY),
    ("CPC",
     "restores revenue by distorting allocation",
     "rewards attention and creates winner-take-share: 66.7% FoodRush under CPC",
     RED),
    ("CPFI",
     "is the aligned mechanism",
     "rewards fulfilled consumer intent · matches neutral welfare · positive platform returns",
     GREEN),
]):
    bx = Inches(0.3) + i * Inches(4.35)
    R(sl, bx, Inches(1.15), Inches(4.1), Inches(4.4), fill=bg)
    T(sl, num, bx, Inches(1.2), Inches(4.1), Inches(1.3),
      sz=52, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    T(sl, headline, bx, Inches(2.52), Inches(4.1), Inches(0.65),
      sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(sl, bx + Inches(0.2), Inches(3.2), Inches(3.7), Pt(1.5), fill=GOLD)
    T(sl, body, bx + Inches(0.1), Inches(3.32), Inches(3.9), Inches(2.0),
      sz=12, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

R(sl, Inches(0.3), Inches(5.78), Inches(12.73), Inches(0.55), fill=NAVY)
T(sl, "Recommendation:  Cost Per Fulfilled Intent  +  bias cap ≤ 0.20",
  Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.42),
  sz=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

T(sl, "The future of AI-mediated marketplaces depends less on whether AI recommends products,\nand more on what the AI is paid to optimize.",
  Inches(0.3), Inches(6.45), Inches(12.73), Inches(0.72),
  sz=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)

notes(sl, "Three takeaways: (1) Neutral LLM is good for consumers but creates negative platform WTP — platforms need a financial reason to participate. (2) CPC restores that reason but rewards the wrong thing — attention rather than satisfaction, producing extreme concentration. (3) CPFI is the right design: it ties the LLM's revenue to fulfilled consumer intent, so the AI has a direct financial incentive to rank by consumer welfare rather than bid size. The policy recommendation is CPFI plus an explicit cap on paid ranking bias at 0.10–0.20. The closing line is the thesis in one sentence.")

