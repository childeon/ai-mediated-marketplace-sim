"""
insert_walkthrough.py
Replaces slide 19 ("One consumer shows the mechanism") in Presentation vF.pptx
with two new slides: A (choice set formation) and B (choice & payoff).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PATH = "Presentation vF.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
TEAL   = RGBColor(0x00, 0x88, 0x8A)
GOLD   = RGBColor(0xF0, 0xAB, 0x00)
GREEN  = RGBColor(0x27, 0x96, 0x58)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xD6, 0xE4, 0xF0)
DARK   = RGBColor(0x22, 0x22, 0x33)
GREY   = RGBColor(0x88, 0x88, 0x99)
SILVER = RGBColor(0xEC, 0xF0, 0xF1)
INK    = RGBColor(0x10, 0x1E, 0x33)
PURPLE = RGBColor(0x6C, 0x35, 0x8A)

SW, SH = Inches(13.33), Inches(7.5)

prs = Presentation(PATH)
BLANK = prs.slide_layouts[0]  # BLANK layout

# ── Helpers ──────────────────────────────────────────────────────────────────

def R(sl, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def T(sl, text, x, y, w, h, sz=12, bold=False, italic=False,
      color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return t

def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text

def footer(sl):
    R(sl, 0, SH - Inches(0.25), SW, Inches(0.25), fill=NAVY)
    T(sl, "Columbia University  ·  Digital Marketplaces  ·  Spring 2026",
      Inches(0.3), SH - Inches(0.23), Inches(12.7), Inches(0.22),
      sz=9, color=RGBColor(0x7F, 0xA7, 0xCC), align=PP_ALIGN.LEFT)

def hdr(sl, text):
    R(sl, 0, 0, SW, Inches(0.95), fill=NAVY)
    R(sl, 0, Inches(0.95), SW, Pt(4), fill=GOLD)
    T(sl, text, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.8),
      sz=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def th(sl, cols, x, y, widths, h=Inches(0.34)):
    cx = x
    for col, w in zip(cols, widths):
        R(sl, cx, y, w, h, fill=NAVY)
        T(sl, col, cx + Inches(0.05), y + Inches(0.03),
          w - Inches(0.1), h - Inches(0.06),
          sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += w

def tr(sl, vals, x, y, widths, bg=WHITE, aligns=None, h=Inches(0.32), sz=9, highlight_col=None, highlight_color=None):
    if aligns is None:
        aligns = [PP_ALIGN.CENTER] * len(vals)
    cx = x
    for i, (val, w, al) in enumerate(zip(vals, widths, aligns)):
        cell_bg = highlight_color if (highlight_col is not None and i == highlight_col) else bg
        R(sl, cx, y, w, h, fill=cell_bg, line=GREY, lw=Pt(0.5))
        T(sl, val, cx + Inches(0.05), y + Inches(0.03),
          w - Inches(0.1), h - Inches(0.06),
          sz=sz, color=DARK if cell_bg != NAVY else WHITE, align=al)
        cx += w

def step_badge(sl, num, text, x, y, color, w=Inches(5.6)):
    R(sl, x, y, Inches(0.26), Inches(0.26), fill=color)
    T(sl, str(num), x, y, Inches(0.26), Inches(0.26),
      sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, text, x + Inches(0.31), y, w - Inches(0.31), Inches(0.26),
      sz=10, color=DARK, align=PP_ALIGN.LEFT)

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_el = slides[old_index]
    xml_slides.remove(slide_el)
    xml_slides.insert(new_index, slide_el)

def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_el = slides[index]
    rId = slide_el.get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_el)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE A — Choice Set Formation
# ══════════════════════════════════════════════════════════════════════════════

sl_a = prs.slides.add_slide(BLANK)
R(sl_a, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl_a)
hdr(sl_a, "One Consumer, Two Discovery Worlds")

# ── Consumer intent strip ────────────────────────────────────────────────────
R(sl_a, 0, Inches(0.99), SW, Inches(0.52), fill=NAVY)
T(sl_a, "Consumer intent:", Inches(0.3), Inches(1.02), Inches(1.5), Inches(0.44),
  sz=9, italic=True, color=LIGHT, align=PP_ALIGN.LEFT)
for i, (label, val, color) in enumerate([
    ("Cuisine", "Korean", GOLD),
    ("Budget", "$50.14", WHITE),
    ("Max time", "45 min", WHITE),
    ("Loyal to", "FoodRush + NomNom", LIGHT),
]):
    bx = Inches(1.9) + i * Inches(2.8)
    T(sl_a, label + ":", bx, Inches(1.01), Inches(1.1), Inches(0.24),
      sz=8, italic=True, color=GREY, align=PP_ALIGN.LEFT)
    T(sl_a, val, bx + Inches(1.05), Inches(1.01), Inches(1.7), Inches(0.44),
      sz=13, bold=True, color=color, align=PP_ALIGN.LEFT)

# ── Panel backgrounds ────────────────────────────────────────────────────────
LX = Inches(0.28)
RX = Inches(6.98)
PW = Inches(6.1)
PY = Inches(1.6)
PH = Inches(5.38)

R(sl_a, LX, PY, PW, PH, fill=WHITE, line=GREY, lw=Pt(1))
R(sl_a, RX, PY, PW, PH, fill=WHITE, line=BLUE, lw=Pt(2))

# Panel headers
R(sl_a, LX, PY, PW, Inches(0.4), fill=RGBColor(0x55, 0x66, 0x77))
T(sl_a, "No-LLM Path  —  Fragmented Search", LX + Inches(0.1), PY + Inches(0.05),
  PW - Inches(0.2), Inches(0.32), sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

R(sl_a, RX, PY, PW, Inches(0.4), fill=BLUE)
T(sl_a, "LLM Path  —  Cross-Platform Aggregation", RX + Inches(0.1), PY + Inches(0.05),
  PW - Inches(0.2), Inches(0.32), sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── No-LLM steps ─────────────────────────────────────────────────────────────
SY = PY + Inches(0.48)
for i, txt in enumerate([
    "Consumer checks FoodRush + NomNom only",
    "Each platform filters to Korean cuisine",
    "Each platform ranks by its own internal score",
    "Consumer sees up to 10 offers per platform",
]):
    step_badge(sl_a, i+1, txt, LX + Inches(0.12), SY + i * Inches(0.32),
               GREY, w=PW - Inches(0.15))

# No-LLM table
TY_NOLLM = SY + Inches(1.4)
nw = [Inches(1.05), Inches(1.1), Inches(0.82), Inches(0.82), Inches(0.95), Inches(1.12)]
th(sl_a, ["Platform", "Offer", "Price", "Time", "Quality", "Plat. Score"],
   LX + Inches(0.12), TY_NOLLM, nw)

_no_llm_rows = [
    ("FoodRush", "R044_Kor", "$50.25", "42 min", "4.0 ★", "0.302"),
    ("FoodRush", "R024_Kor", "$48.83", "32 min", "4.0 ★", "0.255"),
    ("NomNom",   "R006_Kor", "$51.90", "29 min", "4.5 ★", "0.198"),
    ("NomNom",   "R044_Kor", "$50.04", "51 min", "4.0 ★", "0.183"),
]
for i, row in enumerate(_no_llm_rows):
    bg = SILVER if i % 2 == 0 else WHITE
    hl = RGBColor(0xFF, 0xF3, 0xCD) if row[1] == "R044_Kor" and row[0] == "NomNom" else None
    tr(sl_a, row, LX + Inches(0.12), TY_NOLLM + Inches(0.34) + i * Inches(0.32),
       nw, bg=hl if hl else bg)

# Label: chosen in No-LLM
T(sl_a, "★ NomNom R044_Kor chosen by this consumer (51 min → over time limit)",
  LX + Inches(0.12), TY_NOLLM + Inches(0.34) + 4 * Inches(0.32) + Inches(0.05),
  PW - Inches(0.25), Inches(0.28),
  sz=8, italic=True, color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.LEFT)

# ── LLM steps ────────────────────────────────────────────────────────────────
for i, txt in enumerate([
    "LLM queries all 3 platforms simultaneously",
    "LLM scores all Korean offers by intent fit",
    "LLM returns top-5 shortlist to consumer",
    "Consumer clicks (if trust) then decides to buy",
]):
    step_badge(sl_a, i+1, txt, RX + Inches(0.12), SY + i * Inches(0.32), BLUE, w=PW - Inches(0.15))

# LLM table
lw2 = [Inches(1.1), Inches(1.05), Inches(0.82), Inches(0.82), Inches(0.95), Inches(1.12)]
th(sl_a, ["Offer", "Platform", "Price", "Time", "Quality", "LLM Score"],
   RX + Inches(0.12), TY_NOLLM, lw2)

_llm_rows = [
    ("R025_Kor", "QuickEats", "$50.36", "20 min", "4.5 ★", "0.712"),
    ("R006_Kor", "NomNom",    "$51.90", "29 min", "4.5 ★", "0.691"),
    ("R024_Kor", "QuickEats", "$46.79", "28 min", "4.0 ★", "0.678"),
    ("R024_Kor", "NomNom",    "$46.75", "30 min", "4.0 ★", "0.676"),
    ("R024_Kor", "FoodRush",  "$48.83", "32 min", "4.0 ★", "0.670"),
]
for i, row in enumerate(_llm_rows):
    bg = SILVER if i % 2 == 0 else WHITE
    # Highlight QuickEats rows (not in no-LLM)
    hl = RGBColor(0xE8, 0xF5, 0xE9) if row[1] == "QuickEats" else None
    tr(sl_a, row, RX + Inches(0.12), TY_NOLLM + Inches(0.34) + i * Inches(0.32),
       lw2, bg=hl if hl else bg)

T(sl_a, "★ QuickEats not in this consumer's no-LLM choice set — only reachable via LLM",
  RX + Inches(0.12), TY_NOLLM + Inches(0.34) + 5 * Inches(0.32) + Inches(0.05),
  PW - Inches(0.25), Inches(0.28),
  sz=8, italic=True, color=GREEN, align=PP_ALIGN.LEFT)

# ── Bottom callout ────────────────────────────────────────────────────────────
R(sl_a, Inches(0.28), Inches(7.05), Inches(12.77), Inches(0.18), fill=NAVY)
T(sl_a, "The LLM advantage is not magic: it sees QuickEats, which this consumer never checked directly.",
  Inches(0.45), Inches(7.07), Inches(12.43), Inches(0.16),
  sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(sl_a,
    "In the no-LLM world, the consumer is not irrational. They search two familiar platforms and filter "
    "by cuisine. The limitation is fragmented search — they cannot see QuickEats at all. "
    "In the LLM world, all platform offers enter one cross-platform ranking. The shortlist contains "
    "QuickEats offers that this consumer would never have encountered through direct platform browsing. "
    "Platform scores (no-LLM) are based on quality, commission signal, and promotions. "
    "LLM scores are based on cuisine match, budget headroom, time headroom, and quality.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE B — Choice and Payoff
# ══════════════════════════════════════════════════════════════════════════════

sl_b = prs.slides.add_slide(BLANK)
R(sl_b, 0, 0, SW, SH, fill=RGBColor(0xF5, 0xF7, 0xFA))
footer(sl_b)
hdr(sl_b, "Ranking Changes The Choice Set; Choice Is Probabilistic")

# ── Choice mechanics ─────────────────────────────────────────────────────────
R(sl_b, Inches(0.28), Inches(1.05), Inches(6.1), Inches(1.55), fill=SILVER, line=GREY, lw=Pt(1))
R(sl_b, Inches(6.95), Inches(1.05), Inches(6.1), Inches(1.55), fill=RGBColor(0xE8, 0xF0, 0xFB), line=BLUE, lw=Pt(1.5))

T(sl_b, "No-LLM  —  One-stage choice",
  Inches(0.38), Inches(1.08), Inches(5.9), Inches(0.32),
  sz=11, bold=True, color=DARK, align=PP_ALIGN.LEFT)

# No-LLM flow
for i, (lbl, bg, fc) in enumerate([
    ("Visible offers + no-purchase option", RGBColor(0xCC, 0xCC, 0xCC), DARK),
    ("Multinomial logit over utilities",    GREY,                        WHITE),
    ("Purchase  or  No Purchase",           DARK,                        WHITE),
]):
    fx = Inches(0.38) + i * Inches(1.88)
    R(sl_b, fx, Inches(1.42), Inches(1.72), Inches(0.5), fill=bg)
    T(sl_b, lbl, fx + Inches(0.04), Inches(1.44), Inches(1.64), Inches(0.44),
      sz=8, bold=(bg==DARK), color=fc, align=PP_ALIGN.CENTER)
    if i < 2:
        T(sl_b, "→", fx + Inches(1.74), Inches(1.53), Inches(0.14), Inches(0.3),
          sz=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

T(sl_b, "LLM  —  Two-stage choice",
  Inches(7.05), Inches(1.08), Inches(5.9), Inches(0.32),
  sz=11, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

# LLM flow (4 boxes)
_llm_flow = [
    ("Top-5 shortlist\n+ no-click option", RGBColor(0x9B, 0xC2, 0xE6)),
    ("Click logit\n(trust-weighted)", BLUE),
    ("Purchase logit\nor abandon", NAVY),
    ("Order\nor none", INK),
]
for i, (lbl, bg) in enumerate(_llm_flow):
    fx = Inches(7.05) + i * Inches(1.44)
    R(sl_b, fx, Inches(1.42), Inches(1.32), Inches(0.52), fill=bg)
    T(sl_b, lbl, fx + Inches(0.04), Inches(1.44), Inches(1.24), Inches(0.46),
      sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        T(sl_b, "→", fx + Inches(1.34), Inches(1.52), Inches(0.1), Inches(0.3),
          sz=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── Outcome table ─────────────────────────────────────────────────────────────
OY = Inches(2.72)
ow = [Inches(1.55), Inches(2.55), Inches(2.25), Inches(1.6), Inches(1.6), Inches(1.5)]
th(sl_b, ["World", "Ranking Logic", "Chosen Offer", "Intent Fulfil.", "Platform Net", "LLM Payment"],
   Inches(0.28), OY, ow, h=Inches(0.36))

_outcomes = [
    ("No-LLM",     "platform score ranking",    "NomNom  ·  51 min",    "0.763", "$8.51",  "$0.00",  SILVER, GREY),
    ("LLM Neutral","consumer-fit ranking",      "NomNom  ·  29 min",    "0.917", "$10.09", "$0.00",  WHITE,  TEAL),
    ("LLM CPC",    "FoodRush paid boost (+0.30)","FoodRush  ·  42 min", "0.866", "$13.33", "$0.80",  SILVER, RED),
    ("LLM CPFI",   "boost × fulfilment score",  "QuickEats  ·  20 min", "0.931", "$10.47", "$0.61",  WHITE,  GREEN),
]
_world_colors = [GREY, TEAL, RED, GREEN]
for i, (world, logic, offer, IF, pnet, llm_pay, bg, wc) in enumerate(_outcomes):
    ry = OY + Inches(0.36) + i * Inches(0.44)
    cx = Inches(0.28)
    vals = [world, logic, offer, IF, pnet, llm_pay]
    for j, (v, w) in enumerate(zip(vals, ow)):
        cell_bg = wc if j == 0 else bg
        fc = WHITE if j == 0 else DARK
        R(sl_b, cx, ry, w, Inches(0.41), fill=cell_bg, line=GREY, lw=Pt(0.5))
        T(sl_b, v, cx + Inches(0.06), ry + Inches(0.06),
          w - Inches(0.12), Inches(0.3), sz=10,
          bold=(j == 0), color=fc,
          align=PP_ALIGN.LEFT if j in [1, 2] else PP_ALIGN.CENTER)
        cx += w

# ── Two calculation boxes ────────────────────────────────────────────────────
BOX_Y = Inches(4.65)
BOX_H = Inches(2.05)

# Left: IF calculation
R(sl_b, Inches(0.28), BOX_Y, Inches(6.1), BOX_H, fill=INK)
R(sl_b, Inches(0.28), BOX_Y, Inches(6.1), Inches(0.34), fill=GREY)
T(sl_b, "How intent fulfilment was calculated  —  No-LLM chosen offer",
  Inches(0.38), BOX_Y + Inches(0.04), Inches(5.9), Inches(0.26),
  sz=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

T(sl_b, "Offer: Korean, $50.04, 51 min, 4.0 ★",
  Inches(0.4), BOX_Y + Inches(0.42), Inches(5.8), Inches(0.26),
  sz=10, color=LIGHT, align=PP_ALIGN.LEFT)

R(sl_b, Inches(0.4), BOX_Y + Inches(0.74), Inches(5.7), Inches(1.18), fill=RGBColor(0x1E, 0x2E, 0x44))
T(sl_b, "IF  =  cuisine × budget × time × quality",
  Inches(0.5), BOX_Y + Inches(0.8), Inches(5.5), Inches(0.28),
  sz=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
T(sl_b, "=  1.0  ×  1.0  ×  0.88  ×  0.867",
  Inches(0.5), BOX_Y + Inches(1.1), Inches(5.5), Inches(0.26),
  sz=11, color=WHITE, align=PP_ALIGN.LEFT)
T(sl_b, "=  0.763",
  Inches(0.5), BOX_Y + Inches(1.38), Inches(5.5), Inches(0.26),
  sz=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# time penalty note
T(sl_b, "(time penalty: 51 min > 45 min limit → 0.88 = 1 − 6/30 × penalty factor)",
  Inches(0.4), BOX_Y + Inches(1.72), Inches(5.7), Inches(0.22),
  sz=8, italic=True, color=GREY, align=PP_ALIGN.LEFT)

# Right: CPC bias calculation
R(sl_b, Inches(6.95), BOX_Y, Inches(6.1), BOX_H, fill=INK)
R(sl_b, Inches(6.95), BOX_Y, Inches(6.1), Inches(0.34), fill=RED)
T(sl_b, "How FoodRush gets its CPC boost",
  Inches(7.05), BOX_Y + Inches(0.04), Inches(5.9), Inches(0.26),
  sz=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

T(sl_b, "λ = sponsorship bias strength = 0.30 (default)",
  Inches(7.07), BOX_Y + Inches(0.42), Inches(5.8), Inches(0.26),
  sz=10, color=LIGHT, align=PP_ALIGN.LEFT)

R(sl_b, Inches(7.07), BOX_Y + Inches(0.74), Inches(5.7), Inches(1.18), fill=RGBColor(0x3A, 0x10, 0x0A))
T(sl_b, "CPC bias  =  λ × (platform_rate / max_rate)",
  Inches(7.17), BOX_Y + Inches(0.8), Inches(5.5), Inches(0.28),
  sz=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
T(sl_b, "FoodRush  =  0.30 × (0.80 / 0.80)",
  Inches(7.17), BOX_Y + Inches(1.1), Inches(5.5), Inches(0.26),
  sz=11, color=WHITE, align=PP_ALIGN.LEFT)
T(sl_b, "=  0.300  (maximum possible boost)",
  Inches(7.17), BOX_Y + Inches(1.38), Inches(5.5), Inches(0.26),
  sz=14, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

T(sl_b, "(CPFI: boost × predicted IF → FoodRush gets 0.30 × 0.866 = 0.260 instead)",
  Inches(7.07), BOX_Y + Inches(1.72), Inches(5.7), Inches(0.22),
  sz=8, italic=True, color=GREY, align=PP_ALIGN.LEFT)

notes(sl_b,
    "Ranking does not directly enter consumer utility. It changes what the consumer sees. "
    "Then the consumer chooses probabilistically based on cuisine, quality, price, time, promo, "
    "loyalty, trust, and random taste.\n\n"
    "The no-LLM consumer has a one-stage choice. They see offers from their two most-loyal platforms "
    "and choose among those offers plus the outside option.\n\n"
    "The LLM consumer has a two-stage choice. First they click or ignore the shortlist. "
    "If they click, they then either purchase or abandon.\n\n"
    "The paid regimes affect outcomes by changing the LLM shortlist. CPC gives FoodRush a large "
    "additive ranking boost because it has the highest nominal CPC rate. CPFI also uses sponsorship, "
    "but multiplies the boost by predicted intent fulfilment, so poor-fit sponsored offers "
    "receive less amplification.")


# ══════════════════════════════════════════════════════════════════════════════
# REORDER: delete old slide 19, move new slides into position 19
# ══════════════════════════════════════════════════════════════════════════════

total = len(prs.slides)
# New slides are at indices total-2 and total-1
slide_a_idx = total - 2
slide_b_idx = total - 1

# Delete old slide 19 (index 18)
delete_slide(prs, 18)

# Now total-1 and total-2 (slide_b is now at total-2, slide_a at total-3)
# After deletion, indices shift: slide_a is now at (total-3), slide_b at (total-2)
new_total = len(prs.slides)
slide_a_idx = new_total - 2
slide_b_idx = new_total - 1

# Move slide_a to position 18
move_slide(prs, slide_a_idx, 18)
# slide_b is now at new_total-1, move to position 19
move_slide(prs, new_total - 1, 19)

prs.save(PATH)
print(f"Saved. Total slides: {len(prs.slides)}")

# Verify positions
prs2 = Presentation(PATH)
for i in [17, 18, 19, 20]:
    texts = [s.text.strip()[:50] for s in prs2.slides[i].shapes
             if hasattr(s,'text') and s.text.strip() and 'Columbia' not in s.text]
    print(f"Slide {i+1}: {texts[:2]}")
